
# AI-Detect: Web-Based Application (Flask version)
# App built using Python, Flask, and Bootstrap for the web interface

from flask import Flask, render_template, request
import os
import math
import nltk
import torch
import numpy as np
from collections import Counter
from transformers import AutoTokenizer, AutoModelForCausalLM
import io
import docx
import pdfplumber
from pptx import Presentation

# Ensure nltk punkt is available
nltk.download('punkt')

# Load DistilGPT2 model
tokenizer = AutoTokenizer.from_pretrained("distilgpt2")
model = AutoModelForCausalLM.from_pretrained("distilgpt2")
model.eval()

app = Flask(__name__)

def extract_text_from_file(file_storage):
    filename = file_storage.filename.lower()
    content = ""

    if filename.endswith(".txt"):
        content = file_storage.read().decode('utf-8', errors='ignore')

    elif filename.endswith(".docx"):
        doc = docx.Document(file_storage)
        content = "\n".join([p.text for p in doc.paragraphs])

    elif filename.endswith(".pdf"):
        with pdfplumber.open(io.BytesIO(file_storage.read())) as pdf:
            content = "\n".join(page.extract_text() or "" for page in pdf.pages)

    elif filename.endswith(".pptx"):
        prs = Presentation(file_storage)
        text_runs = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_runs.append(shape.text)
        content = "\n".join(text_runs)

    else:
        content = ""  # unsupported file types are ignored silently

    return content.strip()

@app.route('/', methods=['GET', 'POST'])
def index():
    result = None
    if request.method == 'POST':
        text = request.form.get('input_text', '').strip()
        if not text:
            file = request.files.get('input_file')
            if file:
                text = extract_text_from_file(file)

        if not text:
            result = {'error': 'Please enter some text or upload a valid document.'}
        else:
            result = ai_detection_analysis(text)

    return render_template('index.html', result=result)

def calculate_perplexity(text):
    inputs = tokenizer(text, return_tensors='pt', truncation=True, max_length=512)
    input_ids = inputs["input_ids"]
    with torch.no_grad():
        outputs = model(input_ids, labels=input_ids)
        loss = outputs.loss.item()
    return math.exp(loss)

def measure_burstiness(text):
    sentences = nltk.sent_tokenize(text)
    lengths = [len(sentence.split()) for sentence in sentences]
    return round(np.std(lengths) / np.mean(lengths), 2) if lengths else 0

def measure_entropy(text):
    text = text.replace(" ", "")
    freq = Counter(text)
    total = sum(freq.values())
    entropy = -sum((f/total) * math.log2(f/total) for f in freq.values())
    return round(entropy, 2)

def ai_detection_analysis(text):
    result = {}
    try:
        result['perplexity'] = round(calculate_perplexity(text), 2)
    except:
        result['perplexity'] = None

    result['burstiness'] = measure_burstiness(text)
    result['entropy'] = measure_entropy(text)

    score = 0
    checks = 0

    if result['perplexity'] is not None:
        checks += 1
        if result['perplexity'] < 50:
            score += 1

    checks += 1
    if result['burstiness'] < 0.5:
        score += 1

    checks += 1
    if 3.5 <= result['entropy'] <= 4.5:
        score += 1

    result['ai_likelihood_percent'] = round((score / checks) * 100, 1) if checks > 0 else 0
    result['human_likeness'] = round(100 - result['ai_likelihood_percent'], 1)

    return result

if __name__ == '__main__':
    app.run(debug=True)
