
from flask import Flask, render_template, request
from werkzeug.utils import secure_filename
import os
import docx
import pdfplumber
import PyPDF2
import pptx
import re
import math

app = Flask(__name__)

UPLOAD_FOLDER = 'uploads'
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER
ALLOWED_EXTENSIONS = {'txt', 'doc', 'docx', 'pdf', 'ppt', 'pptx'}

def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

def clean_text(text):
    text = re.sub(r'\s+', ' ', text)
    text = re.sub(r'[“”]', '"', text)
    text = re.sub(r"[‘’]", "'", text)
    text = re.sub(r"[^\x00-\x7F]+", "", text)
    return text.strip()

def extract_text_from_file(file):
    filename = secure_filename(file.filename)
    filepath = os.path.join(app.config['UPLOAD_FOLDER'], filename)
    os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)
    file.save(filepath)
    ext = filename.rsplit('.', 1)[1].lower()
    text = ""
    try:
        if ext == 'txt':
            with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
                text = f.read()
        elif (ext == 'docx' or ext == 'doc'):
            doc = docx.Document(filepath)
            text = '\n'.join([para.text for para in doc.paragraphs])
        elif ext == 'pdf':
            with pdfplumber.open(filepath) as pdf:
                text = ''.join(page.extract_text() or '' for page in pdf.pages)
        elif ext == 'pptx':
            presentation = pptx.Presentation(filepath)
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + '\n'
      #  elif ext == 'doc':
      #      text = textract.process(filepath).decode('utf-8')
    except Exception as e:
        print(f"Error extracting text: {e}")
    return clean_text(text)

def calculate_perplexity(text):
    tokens = text.split()
    if not tokens:
        return 0.0
    log_prob = -sum(math.log2((1 / len(tokens))) for _ in tokens)
    return pow(2, log_prob / len(tokens))

def calculate_burstiness(text):
    sentences = re.split(r'[.!?]', text)
    sentence_lengths = [len(sentence.split()) for sentence in sentences if sentence.strip()]
    if len(sentence_lengths) < 2:
        return 0.0
    mean_length = sum(sentence_lengths) / len(sentence_lengths)
    variance = sum((x - mean_length) ** 2 for x in sentence_lengths) / (len(sentence_lengths) - 1)
    return math.sqrt(variance)

def calculate_entropy(text):
    frequency = {}
    for char in text:
        frequency[char] = frequency.get(char, 0) + 1
    total_chars = len(text)
    entropy = -sum((freq / total_chars) * math.log2(freq / total_chars) for freq in frequency.values())
    return entropy

def ai_detection_analysis(text):
    perplexity = calculate_perplexity(text)
    burstiness = calculate_burstiness(text)
    entropy = calculate_entropy(text)
    ai_likelihood = max(0, min(1, (perplexity - 5) / 20 + (burstiness - 2) / 10 - (entropy - 4) / 5))
    human_likeness = 1 - ai_likelihood
    return {
        'perplexity': perplexity,
        'burstiness': burstiness,
        'entropy': entropy,
        'ai_likelihood_percent': round(ai_likelihood * 100, 2),
        'human_likeness': round(human_likeness * 100, 2),
    }

@app.route('/', methods=['GET', 'POST'])
def index():
    result = None
    text = ''
    if request.method == 'POST':
        if 'input_text' in request.form:
            text = request.form['input_text']
        elif 'input_file' in request.files:
            file = request.files['input_file']
            if file and allowed_file(file.filename):
                text = extract_text_from_file(file)
        if text.strip():
            result = ai_detection_analysis(text)
    return render_template('index.html', result=result)

if __name__ == '__main__':
    app.run(debug=True)
